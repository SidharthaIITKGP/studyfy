import streamlit as st
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import HumanMessage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import base64
import os

# --- 1. Caching & Parsing Functions ---

@st.cache_resource
def init_llm(api_key, model_name):
    """
    Initializes the Gemini LLM based on the user-provided model name.
    Caches the connection for that specific model name.
    """
    try:
        return ChatGoogleGenerativeAI(
            model=model_name,
            google_api_key=api_key
        )
    except Exception as e:
        st.sidebar.error(f"Failed to initialize model '{model_name}'. Check the name and API key.")
        return None

@st.cache_data(show_spinner=False)
def parse_ppt_multimodal(file_bytes):
    """
    Parses a PowerPoint file from bytes, extracting text and images
    from each slide, including grouped and placeholder shapes.
    
    --- UPDATED with WMF/EMF skip logic ---
    """
    pptx_file = io.BytesIO(file_bytes)
    prs = Presentation(pptx_file)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_images = []
        
        def find_shapes(shapes):
            for shape in shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text.append(run.text)
                
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    find_shapes(shape.shapes)
                
                # --- START NEW ROBUST IMAGE LOGIC ---
                try:
                    image_data = None
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_data = shape.image
                    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if hasattr(shape, 'image') and shape.image:
                            image_data = shape.image
                    
                    if image_data:
                        # Check content type BEFORE trying to access the blob.
                        # WMF is 'image/x-wmf'. EMF is 'image/x-emf'.
                        content_type = image_data.content_type.lower()
                        if 'wmf' in content_type or 'emf' in content_type:
                            # Log to console that we skipped it (won't show in UI)
                            print(f"Skipping unsupported image (WMF/EMF) on slide {i+1}")
                        else:
                            # This is a safe image (PNG, JPEG), so add it
                            slide_images.append(image_data.blob)
                            
                except Exception as e:
                    # Catch-all for any other unexpected image parsing error
                    print(f"WARNING: Skipping one image on slide {i+1} due to parser error: {e}")
                    pass
                # --- END NEW ROBUST IMAGE LOGIC ---

        find_shapes(slide.shapes)
        slides_data.append({
            "slide_number": i + 1,
            "text": "\n".join(slide_text).strip(),
            "images": slide_images
        })
        
    return slides_data

# --- 2. Core LLM Generation Function ---

def get_gemini_explanation(llm, slide_data, detail_level, include_images, prev_slide_text=None):
    """
    Generates a detailed explanation for a single slide's data
    based on the user's selected options.
    """
    
    # 1. Define the prompt based on the user's desired detail level
    # --- FIX: Corrected the broken "In-Depth" string ---
    detail_prompts = {
        "Summary": "Provide a concise, high-level summary of the following slide content (max 3 sentences).",
        "Standard": "Explain the key concepts on this slide, defining important terms. Keep it clear and focused.",
        "In-Depth": (
            "You are an expert university professor. Your task is to explain the content of this "
            "PowerPoint slide in great detail for a student preparing for an exam.\n"
            "You MUST combine all information (text and images, if provided) into a single, comprehensive explanation.\n"
            "1.  **Analyze the text** and **define all key terms**.\n"
            "2.  **Analyze any images** (diagrams, charts, graphs) and explain what they show.\n"
            "3.  **Connect the text to the images.** (e.g., 'The text mentions X, which is labeled in the diagram...').\n"
            "4.  **Explain the 'why'** behind the concepts. Why is this important?"
        )
    }
    
    instruction_text = detail_prompts[detail_level]

    # 2. Add previous slide context if provided
    if prev_slide_text:
        instruction_text += (
            f"\n\n--- PREVIOUS SLIDE CONTEXT ---\n"
            f"For context, here is the text from the PREVIOUS slide. Use this to inform your explanation "
            f"of the current slide, but do not explain the previous slide itself.\n"
            f"{prev_slide_text}\n"
            f"--- END PREVIOUS SLIDE CONTEXT ---"
        )
    
    # 3. Add the current slide's text
    instruction_text += f"\n\n--- CURRENT SLIDE CONTENT ---"
    if slide_data["text"]:
        instruction_text += f"\n**Slide Text:**\n{slide_data['text']}"
    else:
        instruction_text += "\n[This slide contains no text.]"

    # 4. Prepare the final message parts (text + images)
    message_parts = [{"type": "text", "text": instruction_text}]
    
    # Add images ONLY if the user has ticked the box
    if include_images and slide_data["images"]:
        message_parts.append({"type": "text", "text": "\n**Slide Images:**\n(See attached images)"})
        for img_bytes in slide_data["images"]:
            try:
                b64_image = base64.b64encode(img_bytes).decode('utf-8')
                # --- FIX: Corrected image__url to image_url ---
                message_parts.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{b64_image}"}
                })
            except Exception as e:
                print(f"Error encoding image for API: {e}")
    
    # 5. Create and invoke the LangChain message
    try:
        human_message = HumanMessage(content=message_parts)
        response = llm.invoke([human_message])
        return response.content
    except Exception as e:
        # --- FIX: Indented these lines ---
        st.sidebar.warning(f"Skipped Slide {slide_data['slide_number']}: {e}", icon="⚠️")
        return None # Return None on failure

# --- 3. Download Formatter ---

def format_explanations_for_download(slides_data, explanations, ppt_name):
    """
    Combines all SLIDE text and any GENERATED explanations into a single string.
    """
    output = [f"STUDY GUIDE FOR: {ppt_name}\n"]
    output.append("=" * 40 + "\n\n")
    
    for i, slide in enumerate(slides_data):
        output.append(f"--- SLIDE {slide['slide_number']} ---\n")
        
        if slide['text']:
            output.append("[Slide Text]\n")
            output.append(slide['text'])
            output.append("\n" + "-" * 20 + "\n")
        
        output.append("[Explanation]\n")
        if i < len(explanations):
            output.append(explanations[i])
        else:
            output.append("[Explanation not yet generated.]")
            
        output.append("\n\n" + "=" * 40 + "\n\n")
        
    return "\n".join(output)

# --- 4. Main Streamlit App ---

def main():
    st.set_page_config(
        page_title="Studyfy (Streaming Version)",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    st.sidebar.title("Studyfy")

    # --- 1. API Key & Model Configuration ---
    
    api_key = st.sidebar.text_input(
            "Enter your Google (Gemini) API Key:", 
            type="password",
            key="api_key_input"
        )

    # Dynamic Model Name Input
    # --- FIX: Added unique key ---
    model_name = st.sidebar.text_input(
        "Enter Model Name:",
        value="gemini-2.5-flash",
        help="e.g., 'gemini-1.5-pro-latest', 'gemini-pro-vision'",
        key="model_name_input"
    )

    # --- 2. File Uploader ---
    uploaded_file = st.sidebar.file_uploader(
        "Upload your PowerPoint", 
        type="pptx",
        help="Upload a .pptx file to begin."
    )
    
    if not uploaded_file:
        st.info("Upload a .pptx file in the sidebar to get started.")
        return

    if not api_key or not model_name:
        st.error("Please enter your API Key and a Model Name in the sidebar.")
        return
        
    # --- 3. Analysis Options ---
    st.sidebar.header("Analysis Options")
    
    detail_level = st.sidebar.radio(
        "Level of Detail",
        ["Summary", "Standard", "In-Depth"],
        index=2, 
        help="Choose how detailed you want the explanation to be."
    )
    
    include_images = st.sidebar.checkbox(
        "Analyze Images (Multimodal)",
        value=True,
        help="Allow the AI to 'see' and analyze the images on the slide."
    )
    
    include_context = st.sidebar.checkbox(
        "Include Previous Slide Context",
        value=True,
        help="Provide text from the previous slide to the AI for better context."
    )

    # --- 4. Start Generation Button ---
    if st.sidebar.button("Start Generation", type="primary"):
        
        # We use a try/except block for the *parsing* step
        try:
            with st.spinner("Parsing presentation..."):
                file_bytes = uploaded_file.getvalue()
                all_slides_data = parse_ppt_multimodal(file_bytes)
                llm = init_llm(api_key, model_name)
            
            if not llm:
                st.error("Model initialization failed. Please check the model name.")
                return

            st.success(f"Successfully parsed {len(all_slides_data)} slides. Generating explanations...")
            st.markdown("---")
            st.title("Generated Explanations")
            
            all_explanations = []

            # --- 5. "Streaming" Logic with Dropdowns ---
            for i, slide_data in enumerate(all_slides_data):
                slide_num = slide_data["slide_number"]
                slide_text = slide_data["text"]
                slide_images = slide_data["images"]
                
                # Create one dropdown for *everything* related to this slide
                with st.expander(f"**Slide {slide_num}** - Click to expand/collapse"):

                    # Inner expander for the raw content
                    with st.expander("View Raw Slide Content"):
                        st.markdown("**Extracted Text:**")
                        if not slide_text.strip():
                            st.write("This slide contains no text.")
                        else:
                            st.text(slide_text)
                        
                        st.markdown("**Extracted Images:**")
                        if not slide_images:
                            st.write("This slide contains no images.")
                        else:
                            cols = st.columns(min(len(slide_images), 3))
                            for j, img_bytes in enumerate(slide_images):
                                cols[j % 3].image(img_bytes, use_container_width=True)

                    # --- FIX: This is the robust "skip-on-error" logic ---
                    st.markdown("---")
                    st.header(f"💡 {detail_level} Explanation")
                    
                    explanation = None # Default to None
                    if not slide_text.strip() and not slide_images:
                        explanation = "*(This slide is empty, so no explanation was generated.)*"
                    else:
                        with st.spinner(f"Generating explanation for slide {slide_num}..."):
                            prev_slide_text = all_slides_data[i-1]["text"] if (i > 0 and include_context) else None
                            
                            # This is the call that might return None
                            explanation = get_gemini_explanation(
                                llm,
                                slide_data,
                                detail_level,
                                include_images,
                                prev_slide_text
                            )

                    # Now, check what we got
                    if explanation is None:
                        # The function failed and returned None
                        explanation_to_show = "⚠️ **This slide was skipped due to an error.** (See sidebar for details)"
                        explanation_for_download = "*(Slide skipped due to processing error)*"
                    else:
                        # We have a valid explanation (or the "empty slide" message)
                        explanation_to_show = explanation
                        explanation_for_download = explanation

                    st.markdown(explanation_to_show)
                    all_explanations.append(explanation_for_download) # Save for download

            # --- 6. All processing is done ---
            st.balloons()
            st.success("All slides processed! You can now download the full study guide.")
            
            # --- 7. Add Download Button ---
            file_data = format_explanations_for_download(
                all_slides_data,
                all_explanations,
                uploaded_file.name
            )
            st.sidebar.download_button(
                label="Download Study Guide (.txt)",
                data=file_data,
                file_name=f"{uploaded_file.name.split('.')[0]}_Study_Guide.txt",
                mime="text/plain"
            )

        except Exception as e:
            # This will catch any remaining parsing errors (like the WMF one if it's not caught above)
            st.error(f"A critical error occurred: {e}")
            st.error("This might be a corrupted .pptx file or an issue with the parsing library.")

# Run the app
if __name__ == "__main__":
    main()



